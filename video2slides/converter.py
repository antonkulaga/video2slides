"""Video to PowerPoint converter class."""

import os
import shutil
from datetime import datetime
from pathlib import Path

import cv2
import numpy as np
from eliot import start_action
from pptx import Presentation
from pptx.util import Inches
from skimage.metrics import structural_similarity as ssim


class Video2Slides:
    """Video to Slides converter."""

    def __init__(
        self,
        video_path: str,
        output_path: str | None = None,
        fps_interval: int = 1,
        keep_aspect_ratio: bool = False,
        similarity_threshold: float = 0.95,
        ignore_corners: bool = True,
        corner_size_percent: float = 0.15,
    ) -> None:
        """
        Initialize converter.

        Args:
            video_path: Path to input video file
            output_path: Path to output PPT file
            fps_interval: Extract one frame every N seconds
            keep_aspect_ratio: If True, maintain video aspect ratio in slides
            similarity_threshold: SSIM threshold (0-1) for detecting slide changes (higher = more strict)
            ignore_corners: If True, ignore corner regions when comparing frames (useful for speaker video)
            corner_size_percent: Size of corners to ignore as percentage of frame dimensions (0-1)
        """
        self.video_path = video_path
        self.fps_interval = fps_interval
        self.keep_aspect_ratio = keep_aspect_ratio
        self.similarity_threshold = similarity_threshold
        self.ignore_corners = ignore_corners
        self.corner_size_percent = corner_size_percent
        self.frames: list[str] = []
        self.frames_dir: str | None = None
        self.video_width: int = 0
        self.video_height: int = 0

        if not os.path.exists(video_path):
            raise FileNotFoundError(f"Video file not found: {video_path}")

        if output_path is None:
            base_name = Path(video_path).stem
            output_path = f"{base_name}_slides.pptx"

        # Convert to absolute path
        self.output_path = str(Path(output_path).resolve())

    def _prepare_frame_for_comparison(self, frame: np.ndarray) -> np.ndarray:
        """
        Prepare frame for comparison by optionally masking corners and converting to grayscale.

        Args:
            frame: Input frame in BGR format

        Returns:
            Prepared frame for comparison
        """
        # Convert to grayscale
        gray = cv2.cvtColor(frame, cv2.COLOR_BGR2GRAY)

        if self.ignore_corners:
            # Create a mask to ignore corners where speaker typically appears
            h, w = gray.shape
            mask = np.ones((h, w), dtype=np.uint8) * 255

            # Calculate corner sizes
            corner_h = int(h * self.corner_size_percent)
            corner_w = int(w * self.corner_size_percent)

            # Mask all four corners (typically bottom-right for speaker, but mask all for safety)
            mask[0:corner_h, 0:corner_w] = 0  # Top-left
            mask[0:corner_h, w - corner_w : w] = 0  # Top-right
            mask[h - corner_h : h, 0:corner_w] = 0  # Bottom-left
            mask[h - corner_h : h, w - corner_w : w] = 0  # Bottom-right

            # Apply mask
            gray = cv2.bitwise_and(gray, gray, mask=mask)

        return gray

    def _compute_frame_similarity(self, frame1: np.ndarray, frame2: np.ndarray) -> float:
        """
        Compute similarity between two frames using SSIM.

        Args:
            frame1: First frame (BGR format)
            frame2: Second frame (BGR format)

        Returns:
            Similarity score (0-1, where 1 is identical)
        """
        # Prepare frames for comparison
        gray1 = self._prepare_frame_for_comparison(frame1)
        gray2 = self._prepare_frame_for_comparison(frame2)

        # Resize to reasonable size for faster comparison
        target_height = 480
        if gray1.shape[0] > target_height:
            scale = target_height / gray1.shape[0]
            target_width = int(gray1.shape[1] * scale)
            gray1 = cv2.resize(gray1, (target_width, target_height))
            gray2 = cv2.resize(gray2, (target_width, target_height))

        # Compute SSIM
        similarity_score = ssim(gray1, gray2)

        return similarity_score

    def _is_slide_changed(self, prev_frame: np.ndarray, current_frame: np.ndarray) -> bool:
        """
        Determine if slide content has changed significantly.

        Args:
            prev_frame: Previous frame
            current_frame: Current frame

        Returns:
            True if slide has changed, False otherwise
        """
        similarity = self._compute_frame_similarity(prev_frame, current_frame)
        return similarity < self.similarity_threshold

    def extract_frames(self) -> None:
        """Extract frames from video."""
        with start_action(
            action_type="extract_frames",
            video_path=self.video_path,
            fps_interval=self.fps_interval,
            similarity_threshold=self.similarity_threshold,
            ignore_corners=self.ignore_corners,
        ) as action:
            # Create temporary directory to store frames
            self.frames_dir = "temp_frames"
            os.makedirs(self.frames_dir, exist_ok=True)

            # Open video file
            cap = cv2.VideoCapture(self.video_path)

            if not cap.isOpened():
                raise ValueError(f"Unable to open video file: {self.video_path}")

            fps = cap.get(cv2.CAP_PROP_FPS)
            total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
            self.video_width = int(cap.get(cv2.CAP_PROP_FRAME_WIDTH))
            self.video_height = int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT))
            duration = total_frames / fps if fps > 0 else 0

            action.log(
                message_type="video_info",
                fps=fps,
                total_frames=total_frames,
                duration=duration,
                width=self.video_width,
                height=self.video_height,
            )

            frame_interval = int(fps * self.fps_interval)
            frame_count = 0
            extracted_count = 0
            skipped_count = 0
            prev_frame = None

            while True:
                ret, frame = cap.read()

                if not ret:
                    break

                # Extract frames at specified intervals
                if frame_count % frame_interval == 0:
                    # Check if frame is different from previous
                    should_save = True
                    if prev_frame is not None:
                        should_save = self._is_slide_changed(prev_frame, frame)
                        if not should_save:
                            skipped_count += 1
                            action.log(
                                message_type="frame_skipped",
                                frame_number=frame_count,
                                reason="similar_to_previous",
                            )

                    if should_save:
                        frame_path = os.path.join(
                            self.frames_dir, f"frame_{extracted_count:04d}.jpg"
                        )
                        cv2.imwrite(frame_path, frame)
                        self.frames.append(frame_path)
                        prev_frame = frame.copy()
                        extracted_count += 1

                        if extracted_count % 10 == 0:
                            action.log(
                                message_type="extraction_progress",
                                extracted_count=extracted_count,
                                skipped_count=skipped_count,
                            )

                frame_count += 1

            cap.release()
            action.log(
                message_type="extraction_complete",
                total_extracted=extracted_count,
                total_skipped=skipped_count,
                reduction_ratio=round(skipped_count / (extracted_count + skipped_count) * 100, 2)
                if (extracted_count + skipped_count) > 0
                else 0,
            )

    def generate_ppt(self) -> None:
        """Generate PowerPoint presentation."""
        with start_action(
            action_type="generate_ppt",
            output_path=self.output_path,
            frame_count=len(self.frames),
            keep_aspect_ratio=self.keep_aspect_ratio,
        ) as action:
            if not self.frames:
                raise ValueError("No frame data available")

            # Create presentation
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)

            # Add title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]

            title.text = "Video2Slides"
            subtitle.text = (
                f"Conversion time: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n"
                f"Source file: {os.path.basename(self.video_path)}"
            )

            # Add frame slides
            blank_slide_layout = prs.slide_layouts[6]  # Blank layout

            for idx, frame_path in enumerate(self.frames, 1):
                action.log(message_type="slide_progress", current=idx, total=len(self.frames))

                slide = prs.slides.add_slide(blank_slide_layout)

                if self.keep_aspect_ratio and self.video_width > 0 and self.video_height > 0:
                    # Calculate dimensions maintaining aspect ratio
                    slide_width = prs.slide_width.inches
                    slide_height = prs.slide_height.inches
                    video_aspect = self.video_width / self.video_height
                    slide_aspect = slide_width / slide_height

                    if video_aspect > slide_aspect:
                        # Video is wider than slide
                        width = Inches(slide_width)
                        height = Inches(slide_width / video_aspect)
                        left = Inches(0)
                        top = Inches((slide_height - height.inches) / 2)
                    else:
                        # Video is taller than slide
                        height = Inches(slide_height)
                        width = Inches(slide_height * video_aspect)
                        top = Inches(0)
                        left = Inches((slide_width - width.inches) / 2)
                else:
                    # Fill entire slide (stretch to fit)
                    left = Inches(0)
                    top = Inches(0)
                    width = Inches(10)  # Slide width
                    height = Inches(7.5)  # Slide height

                slide.shapes.add_picture(frame_path, left, top, width=width, height=height)

            # Save presentation
            prs.save(self.output_path)
            action.log(message_type="ppt_saved", output_path=self.output_path)

    def cleanup(self) -> None:
        """Clean up temporary files."""
        with start_action(action_type="cleanup", frames_dir=self.frames_dir):
            if self.frames_dir and os.path.exists(self.frames_dir):
                shutil.rmtree(self.frames_dir)

    def convert(self) -> None:
        """Execute full conversion process."""
        with start_action(
            action_type="convert_video",
            video_path=self.video_path,
            output_path=self.output_path,
            fps_interval=self.fps_interval,
            keep_aspect_ratio=self.keep_aspect_ratio,
        ):
            try:
                self.extract_frames()
                self.generate_ppt()
            finally:
                self.cleanup()
